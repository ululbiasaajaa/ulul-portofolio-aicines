import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Library python-pptx belum terinstall.")
    print("Jalankan: pip install python-pptx")
    sys.exit(1)


# ============================================================
# AICINES-ORIENTED SOFTWARE ENGINEER PORTFOLIO
# ============================================================

OUTPUT_FILENAME = "Portfolio_Muhammad_Ulul_Albab_AICines.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------- Palette ----------
BG = RGBColor(9, 11, 20)
SURFACE = RGBColor(18, 21, 34)
SURFACE_2 = RGBColor(25, 29, 45)
BORDER = RGBColor(48, 53, 76)

WHITE = RGBColor(245, 247, 255)
TEXT = RGBColor(220, 224, 238)
MUTED = RGBColor(145, 151, 174)

PURPLE = RGBColor(139, 92, 246)
VIOLET = RGBColor(168, 85, 247)
CYAN = RGBColor(34, 211, 238)
BLUE = RGBColor(59, 130, 246)
EMERALD = RGBColor(16, 185, 129)
ORANGE = RGBColor(249, 115, 22)
ROSE = RGBColor(244, 63, 94)

FONT = "Aptos"
FONT_BOLD = "Aptos Display"


# ============================================================
# Helpers
# ============================================================

def set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = (
        MSO_SHAPE.ROUNDED_RECTANGLE if radius
        else MSO_SHAPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill

    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line

    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=TEXT,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.0,
):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign

    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align

    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

    return box


def add_rich_text(
    slide,
    runs,
    x,
    y,
    w,
    h,
    size=18,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = align

    for text, color, bold in runs:
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold

    return box


def add_title(slide, eyebrow, title, subtitle=None, accent=PURPLE):
    add_text(
        slide, eyebrow.upper(),
        0.75, 0.42, 11.8, 0.35,
        size=10, color=accent, bold=True
    )
    add_text(
        slide, title,
        0.75, 0.78, 11.8, 0.65,
        size=27, color=WHITE, bold=True,
        font=FONT_BOLD
    )

    if subtitle:
        add_text(
            slide, subtitle,
            0.75, 1.45, 11.6, 0.55,
            size=12.5, color=MUTED
        )

    # small accent line
    add_rect(slide, 0.75, 1.98, 1.15, 0.035, accent, accent, False)


def add_footer(slide, number, total=8):
    add_text(
        slide,
        f"{number:02d} / {total:02d}",
        11.75, 7.03, 0.8, 0.25,
        size=9, color=MUTED, bold=True,
        align=PP_ALIGN.RIGHT
    )


def add_pill(slide, text, x, y, w, color):
    add_rect(slide, x, y, w, 0.36, SURFACE_2, BORDER)
    add_text(
        slide, text,
        x + 0.08, y + 0.04, w - 0.16, 0.25,
        size=9, color=color, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE
    )


def add_bullet_list(slide, bullets, x, y, w, h, size=13, accent=CYAN):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)

    return box


def add_project_card(
    slide,
    x,
    y,
    w,
    h,
    title,
    category,
    description,
    stack,
    accent,
    bullets=None,
    featured=False,
):
    add_rect(slide, x, y, w, h, SURFACE, BORDER)

    # accent rail
    add_rect(slide, x, y, 0.055, h, accent, accent, False)

    if featured:
        add_pill(slide, "FEATURED", x + 0.32, y + 0.28, 0.9, accent)

    title_y = y + (0.78 if featured else 0.32)

    add_text(
        slide, title,
        x + 0.32, title_y, w - 0.65, 0.42,
        size=19 if featured else 16,
        color=WHITE, bold=True,
        font=FONT_BOLD
    )

    add_text(
        slide, category.upper(),
        x + 0.32, title_y + 0.45, w - 0.65, 0.28,
        size=8.5, color=accent, bold=True
    )

    add_text(
        slide, description,
        x + 0.32, title_y + 0.82, w - 0.65, 0.85,
        size=11.5, color=MUTED
    )

    add_text(
        slide, "STACK",
        x + 0.32, title_y + 1.75, 0.6, 0.25,
        size=8.5, color=accent, bold=True
    )

    add_text(
        slide, stack,
        x + 0.32, title_y + 2.02, w - 0.65, 0.55,
        size=10.5, color=TEXT, bold=True
    )

    if bullets:
        bullet_y = title_y + 2.62
        add_bullet_list(
            slide,
            bullets,
            x + 0.32,
            bullet_y,
            w - 0.65,
            h - (bullet_y - y) - 0.25,
            size=10.5 if featured else 10,
            accent=accent,
        )


def add_glow_bar(slide, x, y, w, color):
    # Layered bars create a simple "glow" effect without external assets.
    add_rect(slide, x, y, w, 0.13, color, color, True)
    add_rect(slide, x, y + 0.03, w * 0.72, 0.07, WHITE, WHITE, True)


# ============================================================
# SLIDE 1 — COVER
# ============================================================

slide = prs.slides.add_slide(BLANK)
set_background(slide)

# Decorative blocks
add_rect(slide, 9.45, -0.35, 4.4, 3.0, SURFACE, None)
add_rect(slide, 10.35, 0.2, 3.4, 2.1, SURFACE_2, None)
add_rect(slide, 11.25, 0.75, 2.5, 1.2, PURPLE, None)

add_rect(slide, 0.78, 1.15, 0.07, 4.95, CYAN, CYAN, False)

add_pill(slide, "SOFTWARE ENGINEER", 1.15, 1.18, 1.65, PURPLE)

add_text(
    slide,
    "Muhammad\nUlul Albab",
    1.15, 1.78, 7.5, 1.65,
    size=39, color=WHITE, bold=True,
    font=FONT_BOLD
)

add_text(
    slide,
    "Junior Software Engineer | Full-Stack Web & AI",
    1.15, 3.65, 7.5, 0.5,
    size=19, color=CYAN, bold=True
)

add_text(
    slide,
    "Building practical products across web applications, backend services,\nAI integrations, data systems, and media processing.",
    1.15, 4.25, 7.1, 0.85,
    size=13, color=MUTED
)

add_pill(slide, "React / TypeScript", 1.15, 5.45, 1.55, CYAN)
add_pill(slide, "Node / FastAPI", 2.85, 5.45, 1.45, PURPLE)
add_pill(slide, "AI / RAG", 4.45, 5.45, 1.1, VIOLET)
add_pill(slide, "PostgreSQL", 5.7, 5.45, 1.35, EMERALD)

add_text(
    slide,
    "Bogor, West Java, Indonesia",
    1.15, 6.35, 4.0, 0.3,
    size=10, color=MUTED
)

add_text(
    slide,
    "ulul-portofolio-3odd.vercel.app",
    8.75, 6.35, 3.75, 0.3,
    size=10, color=CYAN, bold=True,
    align=PP_ALIGN.RIGHT
)

add_footer(slide, 1)


# ============================================================
# SLIDE 2 — PROFILE
# ============================================================

slide = prs.slides.add_slide(BLANK)
set_background(slide)
add_title(
    slide,
    "Profile",
    "From Pharmacy to Software Engineering",
    "A practical transition driven by building and shipping real projects.",
    CYAN
)

# Left — profile
add_rect(slide, 0.75, 2.35, 7.15, 3.95, SURFACE, BORDER)

add_text(
    slide,
    "01",
    1.1, 2.7, 0.5, 0.35,
    size=11, color=CYAN, bold=True
)

add_text(
    slide,
    "Hands-on builder",
    1.1, 3.05, 5.9, 0.45,
    size=22, color=WHITE, bold=True,
    font=FONT_BOLD
)

add_text(
    slide,
    "I build full-stack applications and AI-powered products from idea to implementation. "
    "My work spans frontend interfaces, backend APIs, authentication, databases, AI APIs, "
    "vector search, text-to-speech, and cloud deployment.",
    1.1, 3.65, 6.15, 1.25,
    size=13, color=TEXT
)

add_text(
    slide,
    "Current focus",
    1.1, 5.2, 2.0, 0.3,
    size=9, color=CYAN, bold=True
)

add_text(
    slide,
    "Web engineering • AI products • Backend systems",
    1.1, 5.55, 5.9, 0.35,
    size=12, color=WHITE, bold=True
)

# Right — background
add_rect(slide, 8.2, 2.35, 4.35, 3.95, SURFACE_2, BORDER)

add_text(
    slide,
    "02",
    8.55, 2.7, 0.5, 0.35,
    size=11, color=PURPLE, bold=True
)

add_text(
    slide,
    "Analytical foundation",
    8.55, 3.05, 3.5, 0.45,
    size=19, color=WHITE, bold=True,
    font=FONT_BOLD
)

add_text(
    slide,
    "Bachelor's Degree Candidate in Pharmacy\nUniversitas Pakuan • 2023–Present",
    8.55, 3.7, 3.3, 0.75,
    size=12, color=TEXT, bold=True
)

add_text(
    slide,
    "Pharmacy research strengthened my attention to structured data, analytical reasoning, "
    "experimental methods, and systematic problem solving.",
    8.55, 4.75, 3.25, 1.05,
    size=11.5, color=MUTED
)

add_footer(slide, 2)


# ============================================================
# SLIDE 3 — TECH STACK
# ============================================================

slide = prs.slides.add_slide(BLANK)
set_background(slide)
add_title(
    slide,
    "Technical Stack",
    "Tools I Actually Build With",
    "A practical stack across frontend, backend, AI, data, and deployment.",
    PURPLE
)

stack_cards = [
    (
        "Frontend",
        "React\nVite\nTypeScript / JavaScript\nTailwind CSS",
        CYAN,
        0.75,
        2.25,
    ),
    (
        "Backend",
        "Node.js / Express\nFastAPI / Python\nREST APIs\nAuthentication",
        PURPLE,
        3.95,
        2.25,
    ),
    (
        "AI & Media",
        "Gemini API\nRAG / LLM integration\nSentenceTransformers\nTTS / FFmpeg",
        VIOLET,
        7.15,
        2.25,
    ),
    (
        "Data & Cloud",
        "PostgreSQL / pgvector\nFirebase / Supabase\nDocker\nRailway / Vercel",
        EMERALD,
        10.35,
        2.25,
    ),
]

for title, body, accent, x, y in stack_cards:
    add_rect(slide, x, y, 2.25, 3.25, SURFACE, BORDER)
    add_rect(slide, x, y, 2.25, 0.06, accent, accent, False)
    add_text(
        slide, title,
        x + 0.22, y + 0.3, 1.8, 0.4,
        size=15, color=accent, bold=True,
        font=FONT_BOLD
    )
    add_text(
        slide, body,
        x + 0.22, y + 0.95, 1.82, 1.95,
        size=11, color=TEXT
    )

add_text(
    slide,
    "Engineering strengths",
    0.75, 5.85, 2.2, 0.3,
    size=9, color=MUTED, bold=True
)

strengths = [
    ("API Design", CYAN),
    ("Debugging", PURPLE),
    ("Modular Architecture", VIOLET),
    ("Authentication", BLUE),
    ("Git / GitHub", EMERALD),
]

x = 2.45
for label, accent in strengths:
    width = 1.62 if label != "Modular Architecture" else 2.0
    add_pill(slide, label, x, 5.78, width, accent)
    x += width + 0.14

add_footer(slide, 3)


# ============================================================
# SLIDE 4 — PODLEARN
# ============================================================

slide = prs.slides.add_slide(BLANK)
set_background(slide)
add_title(
    slide,
    "Featured Project 01",
    "PodLearn AI",
    "AI podcast and interactive learning platform.",
    VIOLET
)

add_rect(slide, 0.75, 2.25, 4.05, 4.05, SURFACE_2, BORDER)
add_pill(slide, "AI + MEDIA PIPELINE", 1.05, 2.58, 1.65, VIOLET)

add_text(
    slide,
    "From learning material\nto interactive audio.",
    1.05, 3.15, 3.25, 1.05,
    size=23, color=WHITE, bold=True,
    font=FONT_BOLD
)

add_text(
    slide,
    "PDF / TXT",
    1.05, 4.65, 1.05, 0.3,
    size=10, color=CYAN, bold=True
)
add_text(slide, "→", 2.1, 4.61, 0.3, 0.3, size=16, color=MUTED, bold=True)
add_text(
    slide,
    "Gemini",
    2.45, 4.65, 0.9, 0.3,
    size=10, color=PURPLE, bold=True
)
add_text(slide, "→", 3.35, 4.61, 0.3, 0.3, size=16, color=MUTED, bold=True)
add_text(
    slide,
    "TTS",
    3.7, 4.65, 0.65, 0.3,
    size=10, color=VIOLET, bold=True
)

add_glow_bar(slide, 1.05, 5.3, 2.8, VIOLET)

add_text(
    slide,
    "Interactive playback + mid-podcast Q&A",
    1.05, 5.7, 3.25, 0.4,
    size=10.5, color=MUTED
)

add_project_card(
    slide,
    5.1,
    2.25,
    7.45,
    4.05,
    "What I built",
    "End-to-end AI application",
    "A system that converts PDF/text learning materials into structured two-speaker podcast "
    "conversations with auto-generated quizzes.",
    "Node.js • Express • Gemini API • Edge TTS • FFmpeg • Supabase • Google OAuth",
    VIOLET,
    bullets=[
        "Gemini generates structured podcast scripts and quizzes.",
        "Edge TTS produces multi-speaker audio; FFmpeg handles server-side audio processing.",
        "Segment-by-segment playback and dedicated AI Q&A endpoint.",
        "Google OAuth and persistent podcast history through Supabase.",
        "Debugged Railway deployment issues including TTS timeouts and API rate limits.",
    ],
    featured=True,
)

add_footer(slide, 4)


# ============================================================
# SLIDE 5 — SCENTDNA
# ============================================================

slide = prs.slides.add_slide(BLANK)
set_background(slide)
add_title(
    slide,
    "Featured Project 02",
    "ScentDNA",
    "AI semantic search and recommendation engine.",
    CYAN
)

# Architecture visual
add_rect(slide, 0.75, 2.3, 4.0, 4.0, SURFACE, BORDER)

nodes = [
    ("Natural Language", 1.1, 2.75, CYAN),
    ("Embedding", 1.1, 3.55, PURPLE),
    ("pgvector Search", 1.1, 4.35, BLUE),
    ("Gemini RAG", 1.1, 5.15, VIOLET),
]

for label, x, y, accent in nodes:
    add_rect(slide, x, y, 2.75, 0.55, SURFACE_2, accent)
    add_text(
        slide, label,
        x + 0.12, y + 0.12, 2.5, 0.25,
        size=10.5, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER
    )

for y in [3.31, 4.11, 4.91]:
    add_text(slide, "↓", 2.32, y, 0.3, 0.25, size=12, color=MUTED, bold=True)

add_text(
    slide,
    "Semantic retrieval → contextual recommendation",
    1.05, 5.92, 3.3, 0.35,
    size=9.5, color=MUTED
)

add_project_card(
    slide,
    5.05,
    2.3,
    7.5,
    4.0,
    "What I built",
    "AI search + RAG backend",
    "A fragrance search engine that turns natural-language queries into vector embeddings "
    "and retrieves relevant products using PostgreSQL pgvector.",
    "FastAPI • PostgreSQL / pgvector • SentenceTransformers • Gemini • Docker",
    CYAN,
    bullets=[
        "Natural-language queries converted into vector embeddings.",
        "Similarity search performed directly with PostgreSQL pgvector.",
        "Gemini generates contextual recommendations from retrieved results.",
        "FastAPI exposes search and recommendation endpoints with structured schemas.",
        "Optimized model initialization and CPU usage for constrained cloud deployment.",
    ],
    featured=True,
)

add_footer(slide, 5)


# ============================================================
# SLIDE 6 — SCHOOLCOM
# ============================================================

slide = prs.slides.add_slide(BLANK)
set_background(slide)
add_title(
    slide,
    "Featured Project 03",
    "SchoolCom",
    "Role-based school management platform built with React Native and Firebase.",
    BLUE
)

# Big architecture card
add_rect(slide, 0.75, 2.25, 5.1, 4.05, SURFACE, BORDER)

add_text(
    slide,
    "ROLE-BASED PRODUCT",
    1.08, 2.58, 2.0, 0.3,
    size=9, color=BLUE, bold=True
)

roles = [
    ("Admin", "School management", BLUE),
    ("Teacher", "Attendance + assessments", CYAN),
    ("Parent", "Student progress", EMERALD),
]

y = 3.15
for role, desc, accent in roles:
    add_rect(slide, 1.08, y, 1.35, 0.62, SURFACE_2, accent)
    add_text(
        slide, role,
        1.18, y + 0.17, 1.15, 0.25,
        size=10.5, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER
    )
    add_text(
        slide, desc,
        2.68, y + 0.15, 2.35, 0.3,
        size=10.5, color=TEXT
    )
    y += 0.92

add_text(
    slide,
    "Firestore-backed data flows + navigation authorization",
    1.08, 5.95, 4.2, 0.4,
    size=10, color=MUTED
)

add_project_card(
    slide,
    6.1,
    2.25,
    6.45,
    4.05,
    "Engineering focus",
    "Mobile + backend-driven application",
    "An actively developed school management platform covering attendance, assessments, "
    "student progress tracking, and administrative workflows.",
    "React Native • Expo • TypeScript • Firebase / Firestore",
    BLUE,
    bullets=[
        "Implemented Firestore-backed service and UI data flows.",
        "Built role-aware navigation and authorization guards.",
        "Structured the app with modular service and UI layers.",
        "Maintained type safety across an evolving multi-phase roadmap.",
    ],
    featured=True,
)

add_footer(slide, 6)


# ============================================================
# SLIDE 7 — OTHER PROJECTS
# ============================================================

slide = prs.slides.add_slide(BLANK)
set_background(slide)
add_title(
    slide,
    "Additional Projects",
    "More Systems I've Built",
    "Different domains, same engineering mindset: build, debug, ship, improve.",
    ORANGE
)

add_project_card(
    slide,
    0.75,
    2.25,
    5.95,
    1.95,
    "Clinical Suite Dashboard",
    "Web application",
    "Clinical management dashboard with structured data workflows and automated clinical parameter calculators.",
    "HTML5 • JavaScript • Tailwind CSS",
    CYAN,
)

add_project_card(
    slide,
    6.85,
    2.25,
    5.7,
    1.95,
    "LDR Anchor",
    "Mobile application",
    "Android application with real-time mood synchronization, interactive features, and push notifications.",
    "React Native • Firebase Firestore • FCM V1",
    ORANGE,
)

add_project_card(
    slide,
    0.75,
    4.55,
    5.95,
    1.95,
    "Domain Explorer",
    "Backend / web application",
    "Domain management application built with a structured Laravel backend and database workflows.",
    "Laravel • PHP • MySQL",
    EMERALD,
)

add_project_card(
    slide,
    6.85,
    4.55,
    5.7,
    1.95,
    "Cloud Inventory",
    "Serverless application",
    "Inventory workflow with role-based access, audit-oriented operations, and cloud-backed data handling.",
    "Google Apps Script • Google Sheets • Drive",
    PURPLE,
)

add_footer(slide, 7)


# ============================================================
# SLIDE 8 — CLOSING
# ============================================================

slide = prs.slides.add_slide(BLANK)
set_background(slide)

# Decorative right-side visual
add_rect(slide, 8.8, -0.3, 5.0, 8.0, SURFACE, None)
add_rect(slide, 9.55, 0.7, 3.55, 1.2, PURPLE, None)
add_rect(slide, 10.2, 2.2, 2.9, 1.0, CYAN, None)
add_rect(slide, 9.45, 3.7, 3.65, 1.25, VIOLET, None)
add_rect(slide, 10.15, 5.35, 2.95, 0.9, BLUE, None)

add_text(
    slide,
    "LET'S BUILD",
    0.8, 1.2, 3.0, 0.4,
    size=11, color=CYAN, bold=True
)

add_text(
    slide,
    "Ready to contribute,\nlearn fast, and ship.",
    0.8, 1.72, 7.2, 1.45,
    size=35, color=WHITE, bold=True,
    font=FONT_BOLD
)

add_text(
    slide,
    "I'm looking for software engineering opportunities where I can work on "
    "real products, collaborate with engineers, and grow through shipping.",
    0.8, 3.55, 6.5, 1.0,
    size=14, color=TEXT
)

add_pill(slide, "Full-Stack Web", 0.8, 4.95, 1.4, CYAN)
add_pill(slide, "AI Applications", 2.35, 4.95, 1.45, PURPLE)
add_pill(slide, "Backend Systems", 3.95, 4.95, 1.5, EMERALD)

add_text(
    slide,
    "Muhammad Ulul Albab",
    0.8, 5.85, 3.5, 0.35,
    size=14, color=WHITE, bold=True
)

add_text(
    slide,
    "ulula2812@gmail.com\n+62 895-4147-81707\nulul-portofolio-3odd.vercel.app",
    0.8, 6.25, 4.5, 0.8,
    size=10.5, color=MUTED
)

add_text(
    slide,
    "Bogor • Indonesia",
    10.0, 6.7, 2.55, 0.3,
    size=9.5, color=WHITE, bold=True,
    align=PP_ALIGN.RIGHT
)

add_footer(slide, 8)


# ============================================================
# SAVE
# ============================================================

prs.save(OUTPUT_FILENAME)
print(f"SUCCESS: {OUTPUT_FILENAME}")
print(f"Slides: {len(prs.slides)}")
